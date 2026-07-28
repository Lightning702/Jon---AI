from __future__ import annotations

from pathlib import Path
from typing import Any

from app.core.config import DATA_DIR

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.7

THEMES: dict[str, dict[str, str]] = {
    "midnight": {"primary": "1E2761", "secondary": "CADCFC", "accent": "F5B841", "dark": "121A3F", "light": "F7F9FF", "muted": "5B6486", "soft": "E8EDFB"},
    "forest": {"primary": "2C5F2D", "secondary": "97BC62", "accent": "E8B54A", "dark": "1B3A1C", "light": "F4F8F0", "muted": "5A6B52", "soft": "E4EEDA"},
    "coral": {"primary": "E0454B", "secondary": "F9E795", "accent": "2F3C7E", "dark": "331B1D", "light": "FFF6F1", "muted": "8A6560", "soft": "FBE2DE"},
    "terracotta": {"primary": "B85042", "secondary": "E7E8D1", "accent": "6E8B7A", "dark": "3A1E1A", "light": "FBF7F0", "muted": "7A5C52", "soft": "F0E6DA"},
    "ocean": {"primary": "065A82", "secondary": "1C7293", "accent": "F2B950", "dark": "0B2545", "light": "F2F7FA", "muted": "4A6B7D", "soft": "DDEAF2"},
    "charcoal": {"primary": "36454F", "secondary": "8A9BA8", "accent": "E5B53A", "dark": "1C242A", "light": "F4F5F6", "muted": "5E6C75", "soft": "E3E7EA"},
    "teal": {"primary": "028090", "secondary": "00A896", "accent": "F2A65A", "dark": "05313D", "light": "F1FAF8", "muted": "4C7B7C", "soft": "D9EFEA"},
    "berry": {"primary": "6D2E46", "secondary": "A26769", "accent": "D9A566", "dark": "3A1727", "light": "FBF5EE", "muted": "805B62", "soft": "F0E2DC"},
    "sage": {"primary": "50808E", "secondary": "84B59F", "accent": "D9B25F", "dark": "22333B", "light": "F5F8F3", "muted": "5A7068", "soft": "E2EDE3"},
    "cherry": {"primary": "990011", "secondary": "D96C6C", "accent": "2F3C7E", "dark": "3B0007", "light": "FCF6F5", "muted": "7C4B4B", "soft": "F4E0DE"},
    "gold": {"primary": "1A1A1E", "secondary": "5A5A66", "accent": "E5B53A", "dark": "0B0B0D", "light": "F7F5F0", "muted": "5A5A5F", "soft": "EDE8DC"},
}
DEFAULT_THEME = "midnight"
CONTENT_TOP = 2.0
CONTENT_BOTTOM = SLIDE_H - 0.75

FONT_TITLE = "Cambria"
FONT_BODY = "Calibri"


def _theme(name: str) -> dict[str, str]:
    return THEMES.get(str(name).strip().lower(), THEMES[DEFAULT_THEME])


def _rgb(value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(str(value).lstrip("#").upper()[:6])


def _target_path(path: str | None, title: str) -> Path:
    if path:
        target = Path(str(path)).expanduser()
    else:
        safe = "".join(c for c in title if c.isalnum() or c in " -_").strip() or "Praesentation"
        target = DATA_DIR / "praesentationen" / f"{safe}.pptx"
    if target.suffix.lower() != ".pptx":
        target = target.with_suffix(".pptx")
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(slide, color: str):
    from pptx.util import Inches

    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    color: str,
    bold: bool = False,
    font: str = FONT_BODY,
    italic: bool = False,
    align: str = "left",
    anchor: str = "top",
    spacing: float = 0.0,
):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if spacing:
        paragraph.space_after = Pt(spacing)
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    color: str,
    accent: str,
):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(14)
        marker = paragraph.add_run()
        marker.text = "▪  "
        marker.font.size = Pt(size)
        marker.font.color.rgb = _rgb(accent)
        marker.font.name = FONT_BODY
        run = paragraph.add_run()
        text = str(item)
        if ":" in text and len(text.split(":", 1)[0]) <= 28:
            head, tail = text.split(":", 1)
            run.text = head + ":"
            run.font.bold = True
            rest = paragraph.add_run()
            rest.text = tail
            rest.font.size = Pt(size)
            rest.font.name = FONT_BODY
            rest.font.color.rgb = _rgb(color)
        else:
            run.text = text
        run.font.size = Pt(size)
        run.font.name = FONT_BODY
        run.font.color.rgb = _rgb(color)
    return box


def _card(slide, left: float, top: float, width: float, height: float, fill: str, line: str | None = None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _circle(slide, left: float, top: float, size: float, fill: str, label: str, color: str):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.word_wrap = False
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = str(label)
    run.font.size = Pt(int(size * 26))
    run.font.bold = True
    run.font.name = FONT_BODY
    run.font.color.rgb = _rgb(color)
    return shape


def _picture(slide, image: str, left: float, top: float, width: float, height: float) -> bool:
    from pptx.util import Inches

    path = Path(str(image)).expanduser()
    if not path.is_file():
        return False
    try:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    except Exception:
        return False
    return True


def _notes(slide, text: str) -> None:
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = str(text)


def _fit(text: str, big: int, small: int, limit: int) -> int:
    length = len(str(text))
    if length <= limit:
        return big
    if length <= limit * 1.8:
        return int(big - (big - small) * 0.5)
    return small


def _bullet_height(items: list, size: int) -> float:
    return len(items) * (size * 1.4 + 14) / 72


def _centered(needed: float, top: float = CONTENT_TOP, bottom: float = CONTENT_BOTTOM) -> float:
    space = bottom - top
    if needed >= space:
        return top
    return top + (space - needed) * 0.35


def _heading(slide, title: str, colors: dict[str, str], on_dark: bool = False) -> None:
    if not title:
        return
    color = colors["light"] if on_dark else colors["primary"]
    _text(slide, title, MARGIN, 0.62, SLIDE_W - 2 * MARGIN, 1.1, _fit(title, 36, 28, 40), color, bold=True, font=FONT_TITLE)


def _corner(slide, colors: dict[str, str], color: str | None = None) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SLIDE_W - 3.1), Inches(SLIDE_H - 3.1), Inches(5.4), Inches(5.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color or colors["primary"])
    shape.line.fill.background()
    shape.shadow.inherit = False


def _title_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["dark"])
    _corner(slide, colors)
    title = str(data.get("title", ""))
    subtitle = str(data.get("subtitle", data.get("text", "")))
    _text(slide, title, MARGIN, 2.35, SLIDE_W - 2 * MARGIN - 2.6, 2.0, _fit(title, 54, 38, 26), colors["light"], bold=True, font=FONT_TITLE)
    if subtitle:
        _text(slide, subtitle, MARGIN, 4.45, SLIDE_W - 2 * MARGIN - 3.4, 1.0, 20, colors["secondary"], italic=True)
    footer = str(data.get("footer", ""))
    if footer:
        _text(slide, footer, MARGIN, 6.45, SLIDE_W - 2 * MARGIN - 3.4, 0.4, 12, colors["muted"])


def _bullets_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["light"])
    _heading(slide, str(data.get("title", "")), colors)
    items = [str(i) for i in data.get("bullets") or [] if str(i).strip()]
    text = str(data.get("text", ""))
    image = str(data.get("image", ""))
    width = SLIDE_W - 2 * MARGIN
    if image:
        width = 6.9
        _picture(slide, image, MARGIN + width + 0.6, CONTENT_TOP, SLIDE_W - MARGIN - (MARGIN + width + 0.6), 4.2)
    needed = _bullet_height(items, 19) + (1.0 if text else 0)
    top = _centered(needed + 0.9)
    if text:
        _text(slide, text, MARGIN, top, width, 0.9, 19, colors["muted"])
        top += 1.0
    if items:
        _card(slide, MARGIN - 0.4, top - 0.4, width + 0.8, _bullet_height(items, 19) + 0.55, "FFFFFF", colors["soft"])
        _bullets(slide, items, MARGIN, top, width, CONTENT_BOTTOM - top, 19, colors["dark"], colors["primary"])


def _cards_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["light"])
    _heading(slide, str(data.get("title", "")), colors)
    items = list(data.get("items") or data.get("cards") or [])[:4]
    if not items:
        _bullets_slide(prs, slide, data, colors)
        return
    count = len(items)
    gap = 0.45
    width = (SLIDE_W - 2 * MARGIN - gap * (count - 1)) / count
    height = min(CONTENT_BOTTOM - CONTENT_TOP, 3.9)
    top = _centered(height)
    for index, item in enumerate(items):
        left = MARGIN + index * (width + gap)
        _card(slide, left, top, width, height, "FFFFFF", colors["soft"])
        head = str(item.get("title", "") if isinstance(item, dict) else item)
        body = str(item.get("text", "") if isinstance(item, dict) else "")
        _circle(slide, left + 0.4, top + 0.45, 0.7, colors["primary"], str(index + 1), colors["light"])
        _text(slide, head, left + 0.4, top + 1.4, width - 0.8, 0.7, _fit(head, 21, 17, 18), colors["primary"], bold=True, font=FONT_TITLE)
        if body:
            _text(slide, body, left + 0.4, top + 2.15, width - 0.8, height - 2.5, 15, colors["dark"])


def _stat_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["primary"])
    items = list(data.get("items") or data.get("stats") or [])[:3]
    title = str(data.get("title", ""))
    _heading(slide, title, colors, on_dark=True)
    if not items:
        items = [{"title": str(data.get("value", data.get("text", ""))), "text": str(data.get("label", ""))}]
    count = max(len(items), 1)
    width = (SLIDE_W - 2 * MARGIN) / count
    top = _centered(2.6)
    for index, item in enumerate(items):
        left = MARGIN + index * width
        value = str(item.get("title", "") if isinstance(item, dict) else item)
        label = str(item.get("text", "") if isinstance(item, dict) else "")
        _text(slide, value, left, top, width - 0.4, 1.5, _fit(value, 68, 40, 6), colors["accent"], bold=True, font=FONT_TITLE, align="center")
        if label:
            _text(slide, label, left, top + 1.65, width - 0.4, 1.0, 17, colors["light"], align="center")


def _two_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["light"])
    _heading(slide, str(data.get("title", "")), colors)
    width = (SLIDE_W - 2 * MARGIN - 0.6) / 2
    columns = list(data.get("items") or data.get("columns") or [])[:2]
    while len(columns) < 2:
        columns.append({})
    needed = 1.4
    for column in columns:
        if isinstance(column, dict):
            body = column.get("bullets") or []
            needed = max(needed, 1.3 + _bullet_height(list(body), 16) + (0.9 if column.get("text") else 0))
    height = min(CONTENT_BOTTOM - CONTENT_TOP, max(needed, 2.6))
    top = _centered(height)
    for index, column in enumerate(columns):
        left = MARGIN + index * (width + 0.6)
        head = str(column.get("title", "") if isinstance(column, dict) else column)
        body = column.get("bullets") if isinstance(column, dict) else None
        text = str(column.get("text", "") if isinstance(column, dict) else "")
        highlight = index == 1
        _card(slide, left, top, width, height, colors["soft"] if highlight else "FFFFFF", colors["soft"])
        inner = top + 0.45
        if head:
            _text(slide, head, left + 0.45, inner, width - 0.9, 0.6, 22, colors["primary"], bold=True, font=FONT_TITLE)
            inner += 0.85
        if text:
            _text(slide, text, left + 0.45, inner, width - 0.9, 0.9, 16, colors["muted"])
            inner += 0.95
        if body:
            _bullets(slide, [str(b) for b in body], left + 0.45, inner, width - 0.9, top + height - inner - 0.3, 16, colors["dark"], colors["primary"])


def _image_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["light"])
    image = str(data.get("image", ""))
    if not _picture(slide, image, SLIDE_W / 2, 0, SLIDE_W / 2, SLIDE_H):
        _fill_half(slide, colors)
    width = SLIDE_W / 2 - MARGIN - 0.6
    title = str(data.get("title", ""))
    items = [str(i) for i in data.get("bullets") or [] if str(i).strip()]
    text = str(data.get("text", ""))
    needed = 1.5 + (1.3 if text else 0) + _bullet_height(items, 16)
    top = _centered(needed, 1.0, SLIDE_H - 1.0)
    _text(slide, title, MARGIN, top, width, 1.6, _fit(title, 38, 28, 26), colors["primary"], bold=True, font=FONT_TITLE)
    top += 1.5
    if text:
        _text(slide, text, MARGIN, top, width, 1.2, 17, colors["muted"])
        top += 1.3
    if items:
        _bullets(slide, items, MARGIN, top, width, SLIDE_H - top - 0.8, 16, colors["dark"], colors["primary"])


def _fill_half(slide, colors: dict[str, str]) -> None:
    from pptx.util import Inches

    shape = slide.shapes.add_shape(1, Inches(SLIDE_W / 2), Inches(0), Inches(SLIDE_W / 2), Inches(SLIDE_H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(colors["primary"])
    shape.line.fill.background()
    shape.shadow.inherit = False


def _quote_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["dark"])
    quote = str(data.get("text", data.get("title", "")))
    _text(slide, f"„{quote}“", 1.6, 2.1, SLIDE_W - 3.2, 3.0, _fit(quote, 34, 24, 90), colors["light"], italic=True, font=FONT_TITLE, align="center", anchor="middle")
    source = str(data.get("subtitle", data.get("source", "")))
    if source:
        _text(slide, f"— {source}", 1.6, 5.4, SLIDE_W - 3.2, 0.6, 16, colors["accent"], align="center")


def _timeline_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["light"])
    _heading(slide, str(data.get("title", "")), colors)
    items = list(data.get("items") or [])[:5]
    if not items:
        _bullets_slide(prs, slide, data, colors)
        return
    step = min((CONTENT_BOTTOM - CONTENT_TOP) / len(items), 1.35)
    top = _centered(step * len(items))
    for index, item in enumerate(items):
        row = top + index * step
        head = str(item.get("title", "") if isinstance(item, dict) else item)
        body = str(item.get("text", "") if isinstance(item, dict) else "")
        _circle(slide, MARGIN, row, 0.62, colors["primary"], str(index + 1), colors["light"])
        _text(slide, head, MARGIN + 1.0, row, SLIDE_W - MARGIN * 2 - 1.0, 0.45, 20, colors["primary"], bold=True, font=FONT_TITLE)
        if body:
            _text(slide, body, MARGIN + 1.0, row + 0.5, SLIDE_W - MARGIN * 2 - 1.0, step - 0.55, 15, colors["dark"])


def _closing_slide(prs, slide, data: dict, colors: dict[str, str]) -> None:
    _fill(slide, colors["dark"])
    _corner(slide, colors)
    title = str(data.get("title", "Danke!"))
    _text(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN - 2.8, 1.6, _fit(title, 48, 34, 24), colors["light"], bold=True, font=FONT_TITLE)
    subtitle = str(data.get("subtitle", data.get("text", "")))
    if subtitle:
        _text(slide, subtitle, MARGIN, 4.4, SLIDE_W - 2 * MARGIN - 3.4, 1.2, 18, colors["secondary"])


LAYOUTS = {
    "title": _title_slide,
    "bullets": _bullets_slide,
    "cards": _cards_slide,
    "stat": _stat_slide,
    "stats": _stat_slide,
    "two_columns": _two_slide,
    "columns": _two_slide,
    "image": _image_slide,
    "quote": _quote_slide,
    "timeline": _timeline_slide,
    "closing": _closing_slide,
}


class PptxService:
    def themes(self) -> list[str]:
        return sorted(THEMES)

    def create(
        self,
        title: str,
        slides: list[dict[str, Any]],
        path: str | None = None,
        theme: str = DEFAULT_THEME,
        subtitle: str = "",
    ) -> dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return {
                "error": "python-pptx fehlt. Installiere es mit: pip install python-pptx"
            }
        if not slides:
            return {"error": "Keine Folien uebergeben."}
        colors = _theme(theme)
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        made: list[str] = []
        first = str(slides[0].get("layout", "")).lower() if isinstance(slides[0], dict) else ""
        if first != "title":
            slides = [{"layout": "title", "title": title, "subtitle": subtitle}] + list(slides)
        for data in slides[:40]:
            if not isinstance(data, dict):
                data = {"layout": "bullets", "title": str(data)}
            layout = str(data.get("layout", "bullets")).strip().lower()
            builder = LAYOUTS.get(layout, _bullets_slide)
            slide = _blank(prs)
            builder(prs, slide, data, colors)
            _notes(slide, str(data.get("notes", "")))
            made.append(layout if layout in LAYOUTS else "bullets")
        target = _target_path(path, title)
        prs.save(str(target))
        return {
            "ok": True,
            "path": str(target),
            "slides": len(made),
            "layouts": made,
            "theme": _theme_name(theme),
        }

    def read(self, path: str, max_slides: int = 60) -> dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx fehlt. Installiere es mit: pip install python-pptx"}
        source = Path(str(path)).expanduser()
        if not source.is_file():
            return {"error": f"Datei nicht gefunden: {source}"}
        prs = Presentation(str(source))
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(prs.slides, start=1):
            if index > max_slides:
                break
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    value = shape.text_frame.text.strip()
                    if value:
                        texts.append(value)
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides.append({"nummer": index, "text": texts, "notizen": notes})
        return {"path": str(source), "folien": len(slides), "inhalt": slides}


def _theme_name(theme: str) -> str:
    name = str(theme).strip().lower()
    return name if name in THEMES else DEFAULT_THEME


_service: PptxService | None = None


def get_pptx_service() -> PptxService:
    global _service
    if _service is None:
        _service = PptxService()
    return _service
